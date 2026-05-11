"""
AituDesk — Editorial Presentation Builder
==========================================
Генерирует .pptx в редакционном стиле фронтенда AituDesk:
    serif-заголовки, § нумерация секций, тонкие hairlines, реальные скриншоты,
    нативные PowerPoint-переходы и entrance-анимации.

Запуск:
    python build.py
Результат:
    AituDesk-Presentation.pptx
"""

from __future__ import annotations

import copy
import os
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ============================================================================
# DESIGN SYSTEM
# ============================================================================

# Палитра — Dark editorial (#1B1A1F + #F4EEE3 + #E8A246)
PAPER       = RGBColor(0x1B, 0x1A, 0x1F)   # deep charcoal
CARD        = RGBColor(0x26, 0x25, 0x2B)   # элевейтед surface
INK         = RGBColor(0xF4, 0xEE, 0xE3)   # тёплая слоновая кость
INK_MUTED   = RGBColor(0xB8, 0xB0, 0xA4)
INK_DIM     = RGBColor(0x86, 0x80, 0x78)
HAIRLINE    = RGBColor(0x3A, 0x36, 0x40)
RULE_DARK   = RGBColor(0xF4, 0xEE, 0xE3)   # rule на тёмном — светлый

PRIMARY     = RGBColor(0x7A, 0x9D, 0xC4)   # пыльный индиго
ACCENT      = RGBColor(0xE8, 0xA2, 0x46)   # amber
SUCCESS     = RGBColor(0x8B, 0xB5, 0x90)   # mossy sage
WARNING     = RGBColor(0xE8, 0xC0, 0x46)
INFO        = RGBColor(0x9A, 0xB8, 0xD8)

# Типографика — fonts с fallback chain
F_SERIF     = "Georgia"        # Literata fallback (editorial display)
F_SANS      = "Calibri"        # Inter fallback (UI sans)
F_MONO      = "Consolas"       # Geist Mono fallback (numbers, ids)

# Размер слайда — 16:9 wide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Корни путей
ROOT       = Path(__file__).parent.parent
SCREENS    = ROOT / "screenshots"
OUT_FILE   = Path(__file__).parent / "AituDesk-Presentation.pptx"


# ============================================================================
# LOW-LEVEL HELPERS
# ============================================================================

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    """Прямоугольник с настраиваемой заливкой и обводкой."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def add_line(slide, x1, y1, x2, y2, color=HAIRLINE, weight=0.5):
    """Тонкая линейка (hairline)."""
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_text(
    slide, x, y, w, h, text, *,
    font=F_SANS, size=12, bold=False, italic=False,
    color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    tracking=0, line_spacing=1.15,
):
    """Текстовый блок с однородным форматированием."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if tracking:
        # spacing tracking via XML
        rPr = r._r.get_or_add_rPr()
        rPr.set("spc", str(tracking))
    return tb


def add_rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """Много раннов в одном параграфе. runs = [(text, dict-style), ...]."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, style in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = style.get("font", F_SANS)
        f.size = Pt(style.get("size", 12))
        f.bold = style.get("bold", False)
        f.italic = style.get("italic", False)
        f.color.rgb = style.get("color", INK)
        spc = style.get("tracking", 0)
        if spc:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(spc))
    return tb


def add_paragraphs(slide, x, y, w, h, paras, *, anchor=MSO_ANCHOR.TOP):
    """Многопараграфный блок. paras = [(text, dict-style), ...]."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    first = True
    for text, style in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = style.get("align", PP_ALIGN.LEFT)
        p.line_spacing = style.get("line_spacing", 1.25)
        if "space_before" in style:
            p.space_before = Pt(style["space_before"])
        if "space_after" in style:
            p.space_after = Pt(style["space_after"])
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = style.get("font", F_SANS)
        f.size = Pt(style.get("size", 12))
        f.bold = style.get("bold", False)
        f.italic = style.get("italic", False)
        f.color.rgb = style.get("color", INK)
        spc = style.get("tracking", 0)
        if spc:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(spc))
    return tb


def fill_background(slide, color=PAPER):
    """Заливка фона слайда сплошным цветом (через background fill)."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


# ============================================================================
# SHARED CHROME — masthead, page footer, § kicker
# ============================================================================

def add_masthead(slide, page_no: str, total: str = "20"):
    """Верхняя editorial-полоса: лого слева, метаданные справа, hairline снизу."""
    # Brand mark (slab navy + uppercase)
    add_rich(slide, Inches(0.6), Inches(0.32), Inches(4), Inches(0.3), [
        ("Aitu", {"font": F_SERIF, "size": 14, "bold": True, "color": INK}),
        ("DESK ", {"font": F_SERIF, "size": 14, "bold": False, "italic": True, "color": INK}),
        ("  ·  ", {"font": F_SANS, "size": 11, "color": INK_DIM}),
        ("SERVICE DESK", {"font": F_MONO, "size": 9, "color": INK_DIM, "tracking": 200}),
    ])

    # Right side meta
    add_rich(slide, Inches(8.8), Inches(0.32), Inches(4.2), Inches(0.3), [
        ("ВЫПУСК · ", {"font": F_MONO, "size": 9, "color": INK_DIM, "tracking": 200}),
        ("27 АПРЕЛЯ 2026 Г.", {"font": F_MONO, "size": 9, "color": INK, "bold": True, "tracking": 200}),
        ("   ·   ", {"font": F_MONO, "size": 9, "color": INK_DIM}),
        (f"P. {page_no} / {total}", {"font": F_MONO, "size": 9, "color": INK, "bold": True, "tracking": 200}),
    ], align=PP_ALIGN.RIGHT)

    # Hairline под masthead
    add_line(slide, Inches(0.6), Inches(0.78), Inches(12.733), Inches(0.78), color=RULE_DARK, weight=0.75)


def add_footer(slide, label: str, page_no: str, total: str = "20"):
    """Нижняя полоса: подпись секции слева, номер страницы справа."""
    add_line(slide, Inches(0.6), Inches(7.05), Inches(12.733), Inches(7.05), color=HAIRLINE, weight=0.5)
    add_text(slide, Inches(0.6), Inches(7.15), Inches(8), Inches(0.3),
             label, font=F_MONO, size=8.5, color=INK_DIM, tracking=200)
    add_text(slide, Inches(11.7), Inches(7.15), Inches(1.0), Inches(0.3),
             f"{page_no} / {total}", font=F_MONO, size=8.5, color=INK,
             align=PP_ALIGN.RIGHT, tracking=200, bold=True)


def add_kicker(slide, x, y, num: str, label: str, color=ACCENT):
    """Editorial kicker: § XX  CHAPTER NAME — короткая префикс-метка."""
    add_rich(slide, x, y, Inches(8), Inches(0.32), [
        ("§ ", {"font": F_SERIF, "size": 13, "italic": True, "color": color, "bold": True}),
        (f"{num}  ", {"font": F_MONO, "size": 11, "color": color, "bold": True, "tracking": 250}),
        (label.upper(), {"font": F_MONO, "size": 10, "color": INK_MUTED, "tracking": 350}),
    ])


# ============================================================================
# SLIDE TEMPLATES
# ============================================================================

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def slide_cover(prs):
    s = add_blank_slide(prs); fill_background(s)

    # Top masthead row (без hairline снизу — чтобы было дышать)
    add_rich(s, Inches(0.7), Inches(0.42), Inches(6), Inches(0.3), [
        ("Aitu", {"font": F_SERIF, "size": 14, "bold": True, "color": INK}),
        ("DESK", {"font": F_SERIF, "size": 14, "italic": True, "color": INK}),
    ])
    add_text(s, Inches(8.0), Inches(0.42), Inches(4.6), Inches(0.3),
             "ВЫПУСК · ДИПЛОМНАЯ РАБОТА · 2026",
             font=F_MONO, size=9.5, color=INK_DIM, tracking=300, align=PP_ALIGN.RIGHT)

    # Top double rule (editorial signature)
    add_line(s, Inches(0.7), Inches(0.92), Inches(12.633), Inches(0.92), color=INK, weight=1.5)
    add_line(s, Inches(0.7), Inches(0.99), Inches(12.633), Inches(0.99), color=INK, weight=0.5)

    # Hero kicker
    add_text(s, Inches(0.7), Inches(1.25), Inches(8), Inches(0.4),
             "ПРЕЗЕНТАЦИЯ ПРОЕКТА  ·  PROJECT BRIEF",
             font=F_MONO, size=10, color=ACCENT, tracking=400, bold=True)

    # Hero headline — large editorial serif (compact)
    add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(1.85),
             "Service Desk",
             font=F_SERIF, size=110, color=INK, italic=False)
    add_text(s, Inches(0.7), Inches(3.45), Inches(12.0), Inches(1.85),
             "для колледжа.",
             font=F_SERIF, size=110, color=ACCENT, italic=True)

    # Subhead — editorial deck
    add_paragraphs(s, Inches(0.7), Inches(5.50), Inches(8.0), Inches(1.1), [
        ("Полнофункциональная веб-платформа управления IT-заявками",
         {"font": F_SERIF, "size": 17, "italic": True, "color": INK_MUTED, "line_spacing": 1.35}),
        ("с реал-тайм чатом, SLA-контролем, ролевой моделью и базой знаний.",
         {"font": F_SERIF, "size": 17, "italic": True, "color": INK_MUTED, "line_spacing": 1.35}),
    ])

    # Spec sheet (right column)
    spec_x = Inches(9.6); spec_y = Inches(5.10)
    add_text(s, spec_x, spec_y, Inches(3.0), Inches(0.3),
             "СОСТАВ ВЫПУСКА", font=F_MONO, size=8.5, color=INK_DIM, tracking=300, bold=True)
    add_line(s, spec_x, spec_y + Inches(0.32), spec_x + Inches(3.0), spec_y + Inches(0.32), color=INK, weight=0.75)

    items = [
        ("01", "Архитектура и стек"),
        ("02", "Жизненный цикл и SLA"),
        ("03", "11 интерфейсов · скриншоты"),
        ("04", "Realtime, метрики, тесты"),
    ]
    for i, (num, lbl) in enumerate(items):
        ry = spec_y + Inches(0.5 + i * 0.32)
        add_text(s, spec_x, ry, Inches(0.4), Inches(0.28),
                 num, font=F_MONO, size=9.5, color=ACCENT, bold=True, tracking=200)
        add_text(s, spec_x + Inches(0.5), ry, Inches(2.6), Inches(0.28),
                 lbl, font=F_SERIF, size=11.5, italic=True, color=INK)

    # Bottom rule + folio
    add_line(s, Inches(0.7), Inches(7.05), Inches(12.633), Inches(7.05), color=INK, weight=0.5)
    add_text(s, Inches(0.7), Inches(7.15), Inches(8), Inches(0.3),
             "AITUDESK  ·  СЛУЖБА ПОДДЕРЖКИ КОЛЛЕДЖА",
             font=F_MONO, size=8.5, color=INK_DIM, tracking=300)
    add_text(s, Inches(10.5), Inches(7.15), Inches(2.2), Inches(0.3),
             "P. 01 / 20", font=F_MONO, size=8.5, color=INK, bold=True,
             align=PP_ALIGN.RIGHT, tracking=300)

    return s


def slide_section_intro(prs):
    """§ 01 — О проекте: что это, в одной странице."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "02")
    add_kicker(s, Inches(0.6), Inches(1.05), "01", "О проекте  ·  Introduction")

    # Big headline (compact 3-line block)
    add_paragraphs(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(2.4), [
        ("Одна платформа",
         {"font": F_SERIF, "size": 44, "color": INK, "line_spacing": 1.1, "bold": True}),
        ("для всех IT-заявок ",
         {"font": F_SERIF, "size": 44, "color": INK, "italic": True, "line_spacing": 1.1}),
        ("колледжа.",
         {"font": F_SERIF, "size": 44, "color": ACCENT, "italic": True, "line_spacing": 1.1}),
    ])

    # Lede paragraph
    add_paragraphs(s, Inches(0.6), Inches(4.40), Inches(7.4), Inches(2.0), [
        ("AituDesk объединяет роли пользователя, агента и администратора в едином потоке: "
         "заявка создаётся, попадает к нужному агенту по специализации, идёт по SLA-контролю, "
         "решается в чате с реал-тайм обновлениями и закрывается с оценкой удовлетворённости.",
         {"font": F_SERIF, "size": 15, "italic": True, "color": INK_MUTED, "line_spacing": 1.5}),
    ])

    # Side stats column
    box_x = Inches(8.5); box_y = Inches(4.0); col_w = Inches(4.2)
    add_text(s, box_x, box_y, col_w, Inches(0.3),
             "ЦИФРЫ ПРОЕКТА", font=F_MONO, size=8.5, color=INK_DIM, tracking=300, bold=True)
    add_line(s, box_x, box_y + Inches(0.3), box_x + col_w, box_y + Inches(0.3),
             color=INK, weight=0.75)

    rows = [
        ("6", "сервисов в Docker Compose"),
        ("80+", "автоматических тестов"),
        ("11", "ключевых интерфейсов"),
        ("4", "уровня SLA · от 1 ч до 72 ч"),
    ]
    for i, (n, label) in enumerate(rows):
        ry = box_y + Inches(0.50 + i * 0.55)
        add_text(s, box_x, ry, Inches(1.4), Inches(0.55),
                 n, font=F_SERIF, size=28, color=ACCENT, bold=True)
        add_text(s, box_x + Inches(1.5), ry + Inches(0.16), col_w - Inches(1.5), Inches(0.4),
                 label, font=F_SANS, size=11, color=INK_MUTED)

    add_footer(s, "§ 01  ·  О ПРОЕКТЕ", "02")
    return s


def slide_stack(prs):
    """§ 02 — стек технологий, таблица в editorial-стиле."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "03")
    add_kicker(s, Inches(0.6), Inches(1.05), "02", "Технологический стек  ·  The Stack")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.1), [
        ("Modern TypeScript end-to-end.",
         {"font": F_SERIF, "size": 38, "color": INK, "bold": True, "line_spacing": 1.05}),
        ("Один язык от БД-схемы до UI-компонента.",
         {"font": F_SERIF, "size": 19, "italic": True, "color": INK_MUTED, "line_spacing": 1.3, "space_before": 8}),
    ])

    # 4-колонная таблица
    cols = [
        ("01", "Frontend",
         ["React 18 · TS", "Vite 5", "TailwindCSS v3", "shadcn/ui", "Zustand", "Recharts", "Socket.IO client"]),
        ("02", "Backend",
         ["Node 20 · Express 4", "Prisma 5", "Socket.IO", "Zod · JWT · bcrypt", "Multer · Helmet", "Pino logger"]),
        ("03", "Data & Infra",
         ["PostgreSQL 17", "Docker Compose", "Nginx (SPA + WS)", "pgAdmin"]),
        ("04", "Observability",
         ["Prometheus", "Grafana", "Vitest + Supertest", "PDF reports"]),
    ]
    col_w = Inches(2.95); gap = Inches(0.15); start_x = Inches(0.6); top = Inches(4.05)

    for i, (num, title, items) in enumerate(cols):
        cx = start_x + (col_w + gap) * i

        # колонка-карточка
        add_rect(s, cx, top, col_w, Inches(2.85), fill=CARD, line=HAIRLINE, line_w=0.75)
        # верхняя цветная hairline
        add_line(s, cx, top, cx + col_w, top, color=ACCENT if i == 0 else INK, weight=2)

        # № + название
        add_text(s, cx + Inches(0.2), top + Inches(0.18), Inches(0.5), Inches(0.3),
                 num, font=F_MONO, size=10, color=ACCENT, tracking=250, bold=True)
        add_text(s, cx + Inches(0.2), top + Inches(0.45), col_w - Inches(0.4), Inches(0.4),
                 title, font=F_SERIF, size=18, color=INK, italic=True, bold=True)

        # hairline разделитель
        add_line(s, cx + Inches(0.2), top + Inches(0.95),
                 cx + col_w - Inches(0.2), top + Inches(0.95), color=HAIRLINE, weight=0.5)

        # пункты
        for j, item in enumerate(items):
            iy = top + Inches(1.1 + j * 0.24)
            add_text(s, cx + Inches(0.2), iy, Inches(0.15), Inches(0.22),
                     "·", font=F_MONO, size=12, color=ACCENT, bold=True)
            add_text(s, cx + Inches(0.38), iy, col_w - Inches(0.55), Inches(0.22),
                     item, font=F_SANS, size=10.5, color=INK)

    add_footer(s, "§ 02  ·  СТЕК ТЕХНОЛОГИЙ", "03")
    return s


def slide_roles(prs):
    """§ 03 — роли пользователей."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "04")
    add_kicker(s, Inches(0.6), Inches(1.05), "03", "Роли · Permissions")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(1.0), [
        ("Три роли · одна модель доступа.",
         {"font": F_SERIF, "size": 38, "color": INK, "bold": True, "line_spacing": 1.05}),
    ])

    cards = [
        ("01", "USER", "Пользователь",
         ["Создаёт заявки",
          "Видит свои тикеты",
          "Чатится с агентом",
          "Закрывает решённые",
          "Ставит оценку 1–5"], INFO),
        ("02", "AGENT", "Агент",
         ["Видит назначенные",
          "Берёт свободные NEW",
          "Меняет статусы",
          "Публичные + INTERNAL заметки",
          "Учитывает SLA"], SUCCESS),
        ("03", "ADMIN", "Администратор",
         ["Всё, что AGENT",
          "Переназначение",
          "Управление ролями",
          "Редакция базы знаний",
          "Полный дашборд"], ACCENT),
    ]

    col_w = Inches(4.0); gap = Inches(0.13); start_x = Inches(0.6); top = Inches(2.9)

    for i, (num, slug, title, perms, color) in enumerate(cards):
        cx = start_x + (col_w + gap) * i
        add_rect(s, cx, top, col_w, Inches(3.85), fill=CARD, line=HAIRLINE, line_w=0.75)
        add_line(s, cx, top, cx + col_w, top, color=color, weight=2.25)

        add_text(s, cx + Inches(0.3), top + Inches(0.25), Inches(0.6), Inches(0.3),
                 num, font=F_MONO, size=10, color=color, bold=True, tracking=250)

        add_text(s, cx + Inches(0.3), top + Inches(0.55), col_w - Inches(0.6), Inches(0.5),
                 slug, font=F_MONO, size=22, color=INK, bold=True, tracking=200)

        add_text(s, cx + Inches(0.3), top + Inches(1.1), col_w - Inches(0.6), Inches(0.45),
                 title, font=F_SERIF, size=20, italic=True, color=INK_MUTED)

        add_line(s, cx + Inches(0.3), top + Inches(1.7),
                 cx + col_w - Inches(0.3), top + Inches(1.7), color=HAIRLINE, weight=0.5)

        for j, p in enumerate(perms):
            iy = top + Inches(1.85 + j * 0.32)
            add_text(s, cx + Inches(0.3), iy, Inches(0.25), Inches(0.3),
                     "✓", font=F_SANS, size=11, color=color, bold=True)
            add_text(s, cx + Inches(0.55), iy, col_w - Inches(0.85), Inches(0.3),
                     p, font=F_SANS, size=11.5, color=INK)

    add_footer(s, "§ 03  ·  РОЛИ ПОЛЬЗОВАТЕЛЕЙ", "04")
    return s


def slide_lifecycle(prs):
    """§ 04 — жизненный цикл тикета."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "05")
    add_kicker(s, Inches(0.6), Inches(1.05), "04", "Жизненный цикл  ·  Ticket Lifecycle")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(1.0), [
        ("Пять статусов. ",
         {"font": F_SERIF, "size": 38, "color": INK, "bold": True, "line_spacing": 1.05}),
    ])
    add_paragraphs(s, Inches(0.6), Inches(2.55), Inches(12.1), Inches(0.8), [
        ("От «NEW» до «CLOSED» с возможностью REOPENED — каждое состояние понятно агенту и пользователю.",
         {"font": F_SERIF, "size": 15, "italic": True, "color": INK_MUTED, "line_spacing": 1.4}),
    ])

    # Flow chart
    states = [
        ("NEW",          INK,     "Создан, ожидает агента"),
        ("IN_PROGRESS",  INFO,    "В работе · SLA таймер идёт"),
        ("WAITING",      WARNING, "Ждём пользователя · SLA пауза"),
        ("RESOLVED",     SUCCESS, "Решён, ждём подтверждения"),
        ("CLOSED",       INK_DIM, "Закрыт · в архив"),
    ]
    flow_y = Inches(3.7); box_h = Inches(1.4); box_w = Inches(2.2); gap = Inches(0.25)
    start_x = Inches(0.6)

    for i, (label, color, sub) in enumerate(states):
        bx = start_x + (box_w + gap) * i
        add_rect(s, bx, flow_y, box_w, box_h, fill=CARD, line=color, line_w=1.25)
        # цветной chip сверху
        add_rect(s, bx, flow_y, box_w, Inches(0.07), fill=color, line=None)

        add_text(s, bx + Inches(0.18), flow_y + Inches(0.25), box_w - Inches(0.36), Inches(0.4),
                 label, font=F_MONO, size=12.5, color=color, bold=True, tracking=120)
        add_text(s, bx + Inches(0.18), flow_y + Inches(0.7), box_w - Inches(0.36), Inches(0.7),
                 sub, font=F_SANS, size=10, color=INK_MUTED, line_spacing=1.3)

        # стрелка к следующему
        if i < len(states) - 1:
            ax1 = bx + box_w + Inches(0.04)
            ax2 = bx + box_w + gap - Inches(0.04)
            ay = flow_y + box_h / 2
            add_line(s, ax1, ay, ax2, ay, color=INK, weight=1)
            # наконечник
            tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, ax2 - Inches(0.05), ay - Inches(0.05),
                                      Inches(0.1), Inches(0.1))
            tri.shadow.inherit = False
            tri.fill.solid(); tri.fill.fore_color.rgb = INK
            tri.line.fill.background()
            tri.rotation = 0

    # REOPENED back-arrow note
    add_text(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(0.4),
             "↺  REOPENED — возврат из RESOLVED/CLOSED обратно в IN_PROGRESS, SLA пересчитывается.",
             font=F_SANS, size=12, italic=True, color=ACCENT)

    # signature line
    add_paragraphs(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.0), [
        ("Фоновый воркер ",
         {"font": F_SANS, "size": 11.5, "color": INK_MUTED}),
    ])
    add_rich(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.4), [
        ("Фоновый воркер  ", {"font": F_SANS, "size": 11.5, "color": INK_MUTED}),
        ("checkSlaBreaches()", {"font": F_MONO, "size": 11, "color": INK, "bold": True}),
        ("  раз в минуту проставляет  ", {"font": F_SANS, "size": 11.5, "color": INK_MUTED}),
        ("slaBreached=true", {"font": F_MONO, "size": 11, "color": ACCENT, "bold": True}),
        ("  просроченным тикетам.", {"font": F_SANS, "size": 11.5, "color": INK_MUTED}),
    ])

    add_footer(s, "§ 04  ·  ЖИЗНЕННЫЙ ЦИКЛ ТИКЕТА", "05")
    return s


def slide_sla(prs):
    """§ 05 — SLA по приоритетам."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "06")
    add_kicker(s, Inches(0.6), Inches(1.05), "05", "SLA  ·  Service Level Agreements")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.0), [
        ("От 30 минут до 3 суток.",
         {"font": F_SERIF, "size": 38, "color": INK, "bold": True, "line_spacing": 1.05}),
        ("Время реакции и решения зависят от приоритета.",
         {"font": F_SERIF, "size": 16, "italic": True, "color": INK_MUTED, "line_spacing": 1.3, "space_before": 6}),
    ])

    # Таблица SLA
    rows = [
        ("CRITICAL",  ACCENT,  "30 мин",   "1 ч",   "Преподаватель не может вести занятие"),
        ("HIGH",      WARNING, "1 ч",      "4 ч",   "Не работает важная служба, есть workaround"),
        ("MEDIUM",    INFO,    "4 ч",      "24 ч",  "Стандартные запросы пользователей"),
        ("LOW",       SUCCESS, "24 ч",     "72 ч",  "Косметические или плановые задачи"),
    ]
    tx = Inches(0.6); ty = Inches(3.4); table_w = Inches(12.1)
    col_w = [Inches(2.0), Inches(1.7), Inches(1.7), Inches(6.7)]

    # Заголовки
    headers = ["ПРИОРИТЕТ", "РЕАКЦИЯ", "РЕШЕНИЕ", "ОПИСАНИЕ"]
    cx = tx
    for i, h in enumerate(headers):
        add_text(s, cx + Inches(0.15), ty, col_w[i] - Inches(0.3), Inches(0.3),
                 h, font=F_MONO, size=9, color=INK_DIM, tracking=300, bold=True)
        cx += col_w[i]
    add_line(s, tx, ty + Inches(0.32), tx + table_w, ty + Inches(0.32), color=INK, weight=1)

    # Строки
    row_h = Inches(0.78)
    for r_i, (pri, col, react, resolve, desc) in enumerate(rows):
        ry = ty + Inches(0.42) + r_i * row_h
        # цветной chip-индикатор слева
        add_rect(s, tx + Inches(0.15), ry + Inches(0.18), Inches(0.16), Inches(0.4),
                 fill=col, line=None)
        # priority text
        add_text(s, tx + Inches(0.4), ry + Inches(0.22), col_w[0] - Inches(0.55), Inches(0.4),
                 pri, font=F_MONO, size=14, color=INK, bold=True, tracking=120)

        cx = tx + col_w[0]
        add_text(s, cx + Inches(0.15), ry + Inches(0.22), col_w[1] - Inches(0.3), Inches(0.4),
                 react, font=F_SERIF, size=22, italic=True, color=col, bold=True)
        cx += col_w[1]
        add_text(s, cx + Inches(0.15), ry + Inches(0.22), col_w[2] - Inches(0.3), Inches(0.4),
                 resolve, font=F_SERIF, size=22, italic=True, color=INK, bold=True)
        cx += col_w[2]
        add_text(s, cx + Inches(0.15), ry + Inches(0.27), col_w[3] - Inches(0.3), Inches(0.5),
                 desc, font=F_SANS, size=11.5, color=INK_MUTED, line_spacing=1.3)

        # row hairline
        add_line(s, tx, ry + row_h, tx + table_w, ry + row_h, color=HAIRLINE, weight=0.4)

    add_footer(s, "§ 05  ·  SLA ПО ПРИОРИТЕТАМ", "06")
    return s


def slide_screenshot(prs, *, page_no, num, label_short, label_long,
                     headline, lede, bullets, screenshot_file, side="right"):
    """
    Универсальный screenshot showcase: текст с одной стороны, скриншот с другой.
    side="right" — текст слева, картинка справа; side="left" — наоборот.
    """
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, page_no)
    add_kicker(s, Inches(0.6), Inches(1.05), num, f"{label_long}")

    img_path = SCREENS / screenshot_file
    if not img_path.exists():
        raise FileNotFoundError(f"Screenshot not found: {img_path}")

    # текст слева, картинка справа
    if side == "right":
        text_x = Inches(0.6); text_w = Inches(5.2)
        img_x = Inches(6.3);  img_w = Inches(6.5)
    else:
        text_x = Inches(7.5); text_w = Inches(5.2)
        img_x = Inches(0.6);  img_w = Inches(6.5)

    # Headline
    add_paragraphs(s, text_x, Inches(1.65), text_w, Inches(2.0), [
        (headline, {"font": F_SERIF, "size": 32, "color": INK,
                    "bold": True, "line_spacing": 1.05}),
    ])

    # Lede
    add_paragraphs(s, text_x, Inches(3.50), text_w, Inches(1.5), [
        (lede, {"font": F_SERIF, "size": 13.5, "italic": True,
                "color": INK_MUTED, "line_spacing": 1.45}),
    ])

    # Hairline
    add_line(s, text_x, Inches(5.05), text_x + Inches(2.5), Inches(5.05),
             color=ACCENT, weight=1.25)

    # Bullets
    for j, b in enumerate(bullets):
        by = Inches(5.25 + j * 0.32)
        add_text(s, text_x, by, Inches(0.18), Inches(0.3),
                 "·", font=F_MONO, size=13, color=ACCENT, bold=True)
        add_text(s, text_x + Inches(0.22), by, text_w - Inches(0.22), Inches(0.30),
                 b, font=F_SANS, size=11, color=INK)

    # Screenshot frame
    img_y = Inches(1.45)
    img_h = Inches(5.45)

    # subtle drop shadow approximation: чуть тёмнее фона с offset
    shadow = add_rect(s, img_x + Inches(0.06), img_y + Inches(0.06),
                      img_w, img_h, fill=RGBColor(0x10, 0x0F, 0x12), line=None)
    # frame border + image
    add_rect(s, img_x, img_y, img_w, img_h, fill=CARD,
             line=HAIRLINE, line_w=0.75)

    # Insert image, preserving aspect, fitting within frame minus 6pt padding
    pad = Inches(0.06)
    pic = s.shapes.add_picture(str(img_path),
                                img_x + pad, img_y + pad,
                                width=img_w - pad * 2, height=img_h - pad * 2)
    # Картинки могут быть desktop (wide) или mobile (tall) — здесь все desktop ratio
    # pptx сам растянет до заданных width/height; полей хватит — desktop screenshots
    # имеют разные пропорции. Для лучшего вида: контейнер фиксирован, картинка fit-cover-ish.

    add_footer(s, f"§ {num}  ·  {label_short.upper()}", page_no)
    return s


def slide_realtime(prs):
    """§ 16 — Realtime через Socket.IO."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "17")
    add_kicker(s, Inches(0.6), Inches(1.05), "16", "Realtime  ·  Socket.IO")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.0), [
        ("Сообщения · статусы · уведомления —",
         {"font": F_SERIF, "size": 34, "color": INK, "bold": True, "line_spacing": 1.1}),
        ("без перезагрузки страницы.",
         {"font": F_SERIF, "size": 34, "color": ACCENT, "italic": True, "line_spacing": 1.1}),
    ])

    # Events table
    events = [
        ("ticket:join / ticket:leave",  "client → server",  "подписка на комнату тикета"),
        ("typing:start / typing:stop",  "client ↔ server",  "индикатор «печатает…»"),
        ("message:new",                 "server → client",  "новое сообщение в тикете"),
        ("ticket:status",               "server → client",  "смена статуса тикета"),
        ("notification:new",            "server → client",  "персональное уведомление"),
    ]

    tx = Inches(0.6); ty = Inches(3.85); tw = Inches(12.1)
    cw = [Inches(4.6), Inches(3.5), Inches(4.0)]
    headers = ["СОБЫТИЕ", "НАПРАВЛЕНИЕ", "ОПИСАНИЕ"]
    cx = tx
    for i, h in enumerate(headers):
        add_text(s, cx + Inches(0.1), ty, cw[i] - Inches(0.2), Inches(0.3),
                 h, font=F_MONO, size=9, color=INK_DIM, tracking=300, bold=True)
        cx += cw[i]
    add_line(s, tx, ty + Inches(0.32), tx + tw, ty + Inches(0.32), color=INK, weight=1)

    row_h = Inches(0.48)
    for i, (ev, dir_, desc) in enumerate(events):
        ry = ty + Inches(0.42) + i * row_h
        cx = tx
        add_text(s, cx + Inches(0.1), ry + Inches(0.10), cw[0] - Inches(0.2), Inches(0.30),
                 ev, font=F_MONO, size=12, color=INK, bold=True)
        cx += cw[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.12), cw[1] - Inches(0.2), Inches(0.30),
                 dir_, font=F_MONO, size=10, color=PRIMARY, tracking=120)
        cx += cw[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.13), cw[2] - Inches(0.2), Inches(0.30),
                 desc, font=F_SANS, size=11, color=INK_MUTED)

        add_line(s, tx, ry + row_h, tx + tw, ry + row_h, color=HAIRLINE, weight=0.4)

    # JWT note (под таблицей с гарантированным зазором)
    add_rich(s, Inches(0.6), Inches(6.78), Inches(12.1), Inches(0.30), [
        ("Авторизация — ", {"font": F_SANS, "size": 11, "color": INK_MUTED}),
        ("JWT в auth.token при io()", {"font": F_MONO, "size": 11, "color": INK, "bold": True}),
        (" · Nginx прокси с Upgrade: websocket.", {"font": F_SANS, "size": 11, "color": INK_MUTED}),
    ])

    add_footer(s, "§ 16  ·  REALTIME ЧЕРЕЗ SOCKET.IO", "17")
    return s


def slide_monitoring(prs):
    """§ 17 — мониторинг."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "18")
    add_kicker(s, Inches(0.6), Inches(1.05), "17", "Мониторинг  ·  Observability")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.0), [
        ("Prometheus + Grafana,",
         {"font": F_SERIF, "size": 38, "color": INK, "bold": True, "line_spacing": 1.05}),
        ("прямо в Docker Compose.",
         {"font": F_SERIF, "size": 19, "italic": True, "color": INK_MUTED, "line_spacing": 1.3, "space_before": 6}),
    ])

    metrics = [
        ("tickets_total",                       "Counter",   "Всего создано тикетов"),
        ("tickets_resolved_total",              "Counter",   "Всего закрыто тикетов"),
        ("ticket_resolution_duration_seconds",  "Histogram", "Время решения тикета (сек)"),
        ("http_requests_total",                 "Counter",   "HTTP-запросы (method, route, status)"),
        ("http_request_duration_seconds",       "Histogram", "Длительность HTTP-запросов"),
        ("aitudesk_*",                          "Default",   "Node.js metrics: CPU, memory, event loop"),
    ]

    tx = Inches(0.6); ty = Inches(3.7); tw = Inches(12.1)
    cw = [Inches(5.6), Inches(2.0), Inches(4.5)]
    headers = ["МЕТРИКА", "ТИП", "ОПИСАНИЕ"]
    cx = tx
    for i, h in enumerate(headers):
        add_text(s, cx + Inches(0.1), ty, cw[i] - Inches(0.2), Inches(0.3),
                 h, font=F_MONO, size=9, color=INK_DIM, tracking=300, bold=True)
        cx += cw[i]
    add_line(s, tx, ty + Inches(0.32), tx + tw, ty + Inches(0.32), color=INK, weight=1)

    row_h = Inches(0.42)
    for i, (m, t, desc) in enumerate(metrics):
        ry = ty + Inches(0.42) + i * row_h
        cx = tx
        add_text(s, cx + Inches(0.1), ry + Inches(0.08), cw[0] - Inches(0.2), Inches(0.3),
                 m, font=F_MONO, size=11, color=INK, bold=True)
        cx += cw[0]
        add_text(s, cx + Inches(0.1), ry + Inches(0.10), cw[1] - Inches(0.2), Inches(0.3),
                 t, font=F_SANS, size=10, italic=True, color=ACCENT)
        cx += cw[1]
        add_text(s, cx + Inches(0.1), ry + Inches(0.10), cw[2] - Inches(0.2), Inches(0.3),
                 desc, font=F_SANS, size=11, color=INK_MUTED)

        add_line(s, tx, ry + row_h, tx + tw, ry + row_h, color=HAIRLINE, weight=0.4)

    # Endpoint footer (под таблицей с зазором 0.18")
    add_rich(s, Inches(0.6), Inches(6.78), Inches(12.1), Inches(0.30), [
        ("Эндпоинт ", {"font": F_SANS, "size": 11, "color": INK_MUTED}),
        ("GET /metrics", {"font": F_MONO, "size": 11, "color": INK, "bold": True}),
        (" · скрейп Prometheus каждые 15 секунд · Grafana dashboard ID 11159 для Node.js.",
         {"font": F_SANS, "size": 11, "color": INK_MUTED}),
    ])

    add_footer(s, "§ 17  ·  МОНИТОРИНГ", "18")
    return s


def slide_testing(prs):
    """§ 18 — тестирование."""
    s = add_blank_slide(prs); fill_background(s)
    add_masthead(s, "19")
    add_kicker(s, Inches(0.6), Inches(1.05), "18", "Тестирование  ·  Vitest + Supertest")

    add_paragraphs(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.4), [
        ("80+ тестов  ·  ~5 секунд",
         {"font": F_SERIF, "size": 44, "color": INK, "bold": True, "line_spacing": 1.05}),
        ("Без реальной БД — Prisma полностью замокана. Изоляция и повторяемость на любой машине.",
         {"font": F_SERIF, "size": 15, "italic": True, "color": INK_MUTED,
          "line_spacing": 1.4, "space_before": 10}),
    ])

    # 11 модулей в 2 колонках
    modules = [
        ("auth",          11, "register, login, refresh, logout, RBAC"),
        ("tickets",       18, "Создание, фильтры, статусы, рейтинг, доступ"),
        ("users",         10, "Профиль, роли (ADMIN), агенты, RBAC"),
        ("messages",       5, "Получение, INTERNAL запрет для USER"),
        ("dashboard",      8, "Stats по роли, by-day, by-category"),
        ("knowledge",     10, "Поиск, viewCount++, CRUD (ADMIN)"),
        ("notifications",  4, "Список, read-all, read by id"),
        ("reports",        6, "PDF magic bytes, content-type"),
        ("metrics",        8, "Инкремент счётчиков, observe гистограмм"),
        ("sla",            5, "calculateSlaDeadlines для всех приоритетов"),
        ("health",         2, "/api/health, /metrics endpoint"),
    ]

    col_y = Inches(4.30); col_w = Inches(5.95)
    for col_i in range(2):
        cx = Inches(0.6) + col_i * Inches(6.2)
        for j in range(6):
            idx = col_i * 6 + j
            if idx >= len(modules):
                continue
            mod, count, desc = modules[idx]
            ry = col_y + Inches(j * 0.40)

            add_text(s, cx, ry, Inches(0.5), Inches(0.36),
                     str(count), font=F_SERIF, size=18, color=ACCENT,
                     bold=True, italic=True, align=PP_ALIGN.RIGHT)
            add_text(s, cx + Inches(0.6), ry + Inches(0.08), Inches(1.4), Inches(0.32),
                     mod, font=F_MONO, size=11, color=INK, bold=True)
            add_text(s, cx + Inches(2.0), ry + Inches(0.10), col_w - Inches(2.0), Inches(0.3),
                     desc, font=F_SANS, size=10.5, color=INK_MUTED)

    add_footer(s, "§ 18  ·  ТЕСТИРОВАНИЕ", "19")
    return s


def slide_final(prs):
    """Финал — спасибо."""
    s = add_blank_slide(prs); fill_background(s)

    # masthead-style brand top
    add_rich(s, Inches(0.7), Inches(0.42), Inches(6), Inches(0.3), [
        ("Aitu", {"font": F_SERIF, "size": 14, "bold": True, "color": INK}),
        ("DESK", {"font": F_SERIF, "size": 14, "italic": True, "color": INK}),
    ])
    add_text(s, Inches(8.0), Inches(0.42), Inches(4.6), Inches(0.3),
             "КОНЕЦ ВЫПУСКА  ·  THANK YOU",
             font=F_MONO, size=9.5, color=INK_DIM, tracking=300, align=PP_ALIGN.RIGHT)
    add_line(s, Inches(0.7), Inches(0.92), Inches(12.633), Inches(0.92), color=INK, weight=1.5)
    add_line(s, Inches(0.7), Inches(0.99), Inches(12.633), Inches(0.99), color=INK, weight=0.5)

    # Hero — gigantic gratitude (compact)
    add_text(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(2.2),
             "Спасибо",
             font=F_SERIF, size=130, color=INK, bold=True)
    add_text(s, Inches(0.7), Inches(3.80), Inches(12.0), Inches(1.9),
             "за внимание.",
             font=F_SERIF, size=100, color=ACCENT, italic=True)

    # CTA / repo line
    add_line(s, Inches(0.7), Inches(6.2), Inches(4.5), Inches(6.2),
             color=ACCENT, weight=1.25)
    add_text(s, Inches(0.7), Inches(6.32), Inches(8), Inches(0.4),
             "ВОПРОСЫ?  ·  AITUDESK · 2026",
             font=F_MONO, size=11, color=INK, tracking=350, bold=True)

    # bottom rule + folio
    add_line(s, Inches(0.7), Inches(7.05), Inches(12.633), Inches(7.05), color=INK, weight=0.5)
    add_text(s, Inches(0.7), Inches(7.15), Inches(8), Inches(0.3),
             "AITUDESK  ·  СЛУЖБА ПОДДЕРЖКИ КОЛЛЕДЖА",
             font=F_MONO, size=8.5, color=INK_DIM, tracking=300)
    add_text(s, Inches(10.5), Inches(7.15), Inches(2.2), Inches(0.3),
             "P. 20 / 20", font=F_MONO, size=8.5, color=INK, bold=True,
             align=PP_ALIGN.RIGHT, tracking=300)

    return s


# ============================================================================
# ANIMATIONS — slide transitions + entrance effects via raw XML
# ============================================================================

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
P15_NS = "http://schemas.microsoft.com/office/powerpoint/2012/main"


def add_transition(slide, kind: str = "fade", direction: str = "l", speed: str = "med"):
    """
    Добавляет нативный slide transition в XML слайда.
    kind: 'fade' | 'push' | 'cover' | 'split' | 'wipe' | 'morph' | 'zoom'
    direction: для push/cover/wipe — 'l','r','u','d'.
    speed: 'slow' | 'med' | 'fast'.
    """
    sld = slide._element

    # удалить существующий transition если есть
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)

    # transition должен идти после cSld и timing — в нужное место.
    # Согласно schema, порядок: cSld, clrMapOvr?, transition?, timing?
    nsmap_local = {"p": P_NS, "p14": P14_NS}
    transition = etree.SubElement(sld, qn("p:transition"),
                                   nsmap=None,
                                   attrib={"spd": speed})

    if kind == "fade":
        etree.SubElement(transition, qn("p:fade"))
    elif kind == "push":
        etree.SubElement(transition, qn("p:push"), attrib={"dir": direction})
    elif kind == "cover":
        etree.SubElement(transition, qn("p:cover"), attrib={"dir": direction})
    elif kind == "split":
        etree.SubElement(transition, qn("p:split"),
                          attrib={"orient": "horz", "dir": "out"})
    elif kind == "wipe":
        etree.SubElement(transition, qn("p:wipe"), attrib={"dir": direction})
    elif kind == "zoom":
        etree.SubElement(transition, qn("p:zoom"), attrib={"dir": "in"})
    elif kind == "morph":
        # Morph живёт в p159 namespace (PowerPoint 2019+)
        # Нужно регистрировать namespace и использовать extLst
        ext_lst = etree.SubElement(transition, qn("p:extLst"))
        ext = etree.SubElement(ext_lst, qn("p:ext"),
                                attrib={"uri": "{C5C82F0A-F2E2-4B22-840E-3DA68DC0D49E}"})
        morph_el = etree.SubElement(ext, "{http://schemas.microsoft.com/office/powerpoint/2015/main}morph",
                                     nsmap={"p159": "http://schemas.microsoft.com/office/powerpoint/2015/main"},
                                     attrib={"option": "byObject"})
    else:
        raise ValueError(f"Unknown transition kind: {kind}")

    # Гарантируем правильный порядок дочерних элементов в <p:sld>
    _reorder_sld_children(sld)


def _reorder_sld_children(sld_element):
    """В <p:sld> дети должны идти в порядке: cSld, clrMapOvr?, transition?, timing?, extLst?"""
    order = ["cSld", "clrMapOvr", "transition", "timing", "extLst"]
    children = {}
    for tag in order:
        node = sld_element.find(qn(f"p:{tag}"))
        if node is not None:
            children[tag] = node
            sld_element.remove(node)
    for tag in order:
        if tag in children:
            sld_element.append(children[tag])


def add_simple_fade_in_animation(slide, shape_indices: list[int], delay_per: float = 0.15):
    """
    Добавляет fade-in entrance анимацию для перечисленных shape (по индексу в slide.shapes).
    Активируется автоматически при показе слайда (onPrev event = «при появлении слайда»).
    Использует timing tree XML — упрощённый, но валидный.

    Каскадная задержка delay_per секунд между элементами.
    """
    if not shape_indices:
        return

    sld = slide._element

    # удалить старый timing если есть
    for old in sld.findall(qn("p:timing")):
        sld.remove(old)

    timing = etree.SubElement(sld, qn("p:timing"))
    tnLst = etree.SubElement(timing, qn("p:tnLst"))
    par_root = etree.SubElement(tnLst, qn("p:par"))
    cTn_root = etree.SubElement(par_root, qn("p:cTn"), attrib={
        "id": "1", "dur": "indefinite", "restart": "never", "nodeType": "tmRoot",
    })
    childTnLst_root = etree.SubElement(cTn_root, qn("p:childTnLst"))
    seq = etree.SubElement(childTnLst_root, qn("p:seq"), attrib={
        "concurrent": "1", "nextAc": "seek",
    })
    cTn_seq = etree.SubElement(seq, qn("p:cTn"), attrib={
        "id": "2", "dur": "indefinite", "nodeType": "mainSeq",
    })
    childTnLst_seq = etree.SubElement(cTn_seq, qn("p:childTnLst"))

    next_id = 3
    for i, shape_idx in enumerate(shape_indices):
        if shape_idx >= len(slide.shapes):
            continue
        shp = slide.shapes[shape_idx]
        sp_id = shp.shape_id
        delay = int(delay_per * 1000 * i)

        # par: один шаг последовательности
        par_step = etree.SubElement(childTnLst_seq, qn("p:par"))
        cTn_step = etree.SubElement(par_step, qn("p:cTn"), attrib={
            "id": str(next_id), "fill": "hold",
        })
        next_id += 1
        stCondLst = etree.SubElement(cTn_step, qn("p:stCondLst"))
        etree.SubElement(stCondLst, qn("p:cond"), attrib={"delay": "indefinite"})
        childTnLst_step = etree.SubElement(cTn_step, qn("p:childTnLst"))

        # внутренний par с условием задержки
        par_inner = etree.SubElement(childTnLst_step, qn("p:par"))
        cTn_inner = etree.SubElement(par_inner, qn("p:cTn"), attrib={
            "id": str(next_id), "fill": "hold",
        })
        next_id += 1
        stCondLst_inner = etree.SubElement(cTn_inner, qn("p:stCondLst"))
        etree.SubElement(stCondLst_inner, qn("p:cond"), attrib={"delay": str(delay)})
        childTnLst_inner = etree.SubElement(cTn_inner, qn("p:childTnLst"))

        # par для эффекта
        par_eff = etree.SubElement(childTnLst_inner, qn("p:par"))
        cTn_eff = etree.SubElement(par_eff, qn("p:cTn"), attrib={
            "id": str(next_id), "presetID": "10",
            "presetClass": "entr", "presetSubtype": "0",
            "fill": "hold", "grpId": "0", "nodeType": "withEffect",
        })
        next_id += 1
        stCondLst_eff = etree.SubElement(cTn_eff, qn("p:stCondLst"))
        etree.SubElement(stCondLst_eff, qn("p:cond"), attrib={"delay": "0"})
        childTnLst_eff = etree.SubElement(cTn_eff, qn("p:childTnLst"))

        # set effect — делает элемент visible
        set_el = etree.SubElement(childTnLst_eff, qn("p:set"))
        cBhvr_set = etree.SubElement(set_el, qn("p:cBhvr"))
        cTn_set = etree.SubElement(cBhvr_set, qn("p:cTn"), attrib={
            "id": str(next_id), "dur": "1", "fill": "hold",
        })
        next_id += 1
        stCondLst_set = etree.SubElement(cTn_set, qn("p:stCondLst"))
        etree.SubElement(stCondLst_set, qn("p:cond"), attrib={"delay": "0"})
        tgtEl_set = etree.SubElement(cBhvr_set, qn("p:tgtEl"))
        etree.SubElement(tgtEl_set, qn("p:spTgt"), attrib={"spid": str(sp_id)})
        attrNameLst = etree.SubElement(cBhvr_set, qn("p:attrNameLst"))
        attrName = etree.SubElement(attrNameLst, qn("p:attrName"))
        attrName.text = "style.visibility"
        to_el = etree.SubElement(set_el, qn("p:to"))
        strVal = etree.SubElement(to_el, qn("p:strVal"), attrib={"val": "visible"})

        # animEffect — fade
        anim_eff = etree.SubElement(childTnLst_eff, qn("p:animEffect"), attrib={
            "transition": "in", "filter": "fade",
        })
        cBhvr_anim = etree.SubElement(anim_eff, qn("p:cBhvr"))
        cTn_anim = etree.SubElement(cBhvr_anim, qn("p:cTn"), attrib={
            "id": str(next_id), "dur": "500",
        })
        next_id += 1
        tgtEl_anim = etree.SubElement(cBhvr_anim, qn("p:tgtEl"))
        etree.SubElement(tgtEl_anim, qn("p:spTgt"), attrib={"spid": str(sp_id)})

    # prevCondLst / nextCondLst — нажатие пробела/кликом продолжаем
    prevCondLst = etree.SubElement(seq, qn("p:prevCondLst"))
    cond_prev = etree.SubElement(prevCondLst, qn("p:cond"), attrib={"evt": "onPrev", "delay": "0"})
    tgtEl_prev = etree.SubElement(cond_prev, qn("p:tgtEl"))
    etree.SubElement(tgtEl_prev, qn("p:sldTgt"))
    nextCondLst = etree.SubElement(seq, qn("p:nextCondLst"))
    cond_next = etree.SubElement(nextCondLst, qn("p:cond"), attrib={"evt": "onNext", "delay": "0"})
    tgtEl_next = etree.SubElement(cond_next, qn("p:tgtEl"))
    etree.SubElement(tgtEl_next, qn("p:sldTgt"))

    # bldLst — необходим для каждого участника
    bldLst = etree.SubElement(timing, qn("p:bldLst"))
    for shape_idx in shape_indices:
        if shape_idx >= len(slide.shapes):
            continue
        sp_id = slide.shapes[shape_idx].shape_id
        bld = etree.SubElement(bldLst, qn("p:bldP"), attrib={
            "spid": str(sp_id), "grpId": "0",
        })

    _reorder_sld_children(sld)


# ============================================================================
# BUILD
# ============================================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # === 01 Cover ===
    s = slide_cover(prs)
    add_transition(s, "fade", speed="med")

    # === 02 Intro ===
    s = slide_section_intro(prs)
    add_transition(s, "push", direction="l", speed="med")

    # === 03 Stack ===
    s = slide_stack(prs)
    add_transition(s, "push", direction="l", speed="med")

    # === 04 Roles ===
    s = slide_roles(prs)
    add_transition(s, "push", direction="l", speed="med")

    # === 05 Lifecycle ===
    s = slide_lifecycle(prs)
    add_transition(s, "fade", speed="med")

    # === 06 SLA ===
    s = slide_sla(prs)
    add_transition(s, "push", direction="l", speed="med")

    # ============================================================
    # === 07–16 Screenshot showcases ===
    # ============================================================
    screenshots_meta = [
        # (page_no, num, label_short, label_long, headline, lede, bullets, file, side)
        ("07", "06", "Аутентификация", "Аутентификация  ·  Login",
         "Тихий вход в систему.",
         "Editorial-форма с тонкой типографикой, JWT в куки, refresh-flow прозрачно для пользователя.",
         ["Email + пароль · валидация Zod",
          "Access (15 мин) + Refresh (7 дн)",
          "Запоминание сессии",
          "Линк на регистрацию"],
         "01-login_desktop.png", "right"),

        ("08", "07", "Регистрация", "Регистрация  ·  Sign Up",
         "Минимум полей,",
         "максимум скорости. Bcrypt-хэш паролей, проверка уникальности email на стороне БД.",
         ["Имя, email, пароль",
          "Принудительная сложность пароля",
          "По умолчанию роль USER",
          "Сразу автологин"],
         "02-register_desktop.png", "right"),

        ("09", "08", "Дашборд", "Дашборд  ·  Operational Overview",
         "4 KPI · 2 графика · лента.",
         "Editorial-вёрстка с § нумерацией секций. Спарклайны, donut-chart по категориям, лента свежих заявок и топ агентов — всё на одном экране.",
         ["KPI: всего, открытые, решено, SLA",
          "Линейный график за 14 дней",
          "Donut: распределение по категориям",
          "Топ агентов по решённым тикетам",
          "Сводка SLA"],
         "03-dashboard_desktop.png", "right"),

        ("10", "09", "Лента тикетов", "Лента тикетов  ·  Tickets List",
         "Статус-табы, фильтры, поиск.",
         "Серверный поиск с дебаунсом, фильтрация по приоритету и категории, бейджи статуса и SLA.",
         ["Табы NEW · IN_PROGRESS · WAITING · …",
          "Фильтры по приоритету и категории",
          "Серверный поиск (debounce)",
          "SLA-индикатор у каждой строки"],
         "04-tickets_desktop.png", "right"),

        ("11", "10", "Деталь тикета + чат", "Деталь тикета  ·  Detail + Live Chat",
         "Двухколоночный editorial-лейаут.",
         "Слева — детали и описание. Справа — sticky-чат: публичные сообщения и INTERNAL заметки агентов, реал-тайм через Socket.IO с typing-индикатором.",
         ["Sticky-чат с typing «…»",
          "INTERNAL заметки только агентам",
          "Быстрые действия по роли",
          "Модуль оценки 1–5 после CLOSED",
          "Загрузка вложений до 5 файлов"],
         "11-ticket-detail_desktop.png", "right"),

        ("12", "11", "Создание тикета", "Создание тикета  ·  New Ticket",
         "Форма с подсказками",
         "из базы знаний. Дебаунс-поиск по описанию проблемы — пользователь часто решает вопрос ещё до отправки.",
         ["Категория + приоритет (плитки)",
          "Описание + вложения до 5 файлов",
          "Подсказки KB по введённому тексту",
          "Авто-назначение по специализации"],
         "05-tickets-create_desktop.png", "right"),

        ("13", "12", "База знаний", "База знаний  ·  Knowledge Base",
         "Hero-поиск, плитки, статьи.",
         "Editorial-вёрстка статей с Markdown, viewCount автоматически инкрементируется, поиск по заголовку и контенту.",
         ["Hero-поиск с горячими запросами",
          "Плитки разделов · 4 категории",
          "Список статей с превью",
          "Markdown-редактор для ADMIN"],
         "06-kb_desktop.png", "right"),

        ("14", "13", "Уведомления", "Уведомления  ·  Notifications",
         "Лента с иконкой по типу.",
         "NEW_MESSAGE, TICKET_ASSIGNED, STATUS_CHANGED, TICKET_RATED — четыре типа, прямая ссылка на источник.",
         ["Realtime через Socket.IO",
          "«Прочитать всё» одной кнопкой",
          "Бейдж непрочитанных в шапке",
          "Глубокая ссылка на тикет"],
         "09-notifications_desktop.png", "right"),

        ("15", "14", "Админ · пользователи", "Админ-панель  ·  Users Management",
         "Управление ролями",
         "и специализациями агентов. Назначение специализаций определяет авто-распределение новых тикетов.",
         ["Список всех пользователей",
          "Смена роли USER ↔ AGENT ↔ ADMIN",
          "Назначение специализаций агенту",
          "Активация / деактивация аккаунта"],
         "10-users_desktop.png", "right"),

        ("16", "15", "Профиль", "Профиль  ·  Profile",
         "Личные данные и аватар.",
         "Multipart-загрузка изображений через Multer, обновление в реальном времени, сохранение при смене темы.",
         ["Имя, email, аватар",
          "Загрузка PNG / JPEG до 5 МБ",
          "Светлая / тёмная тема",
          "История последних действий"],
         "08-profile_desktop.png", "right"),
    ]

    push_dirs = ["l", "l", "l", "l", "l", "l", "l", "l", "l", "l"]  # стабильное направление
    for i, meta in enumerate(screenshots_meta):
        s = slide_screenshot(prs,
                              page_no=meta[0], num=meta[1],
                              label_short=meta[2], label_long=meta[3],
                              headline=meta[4], lede=meta[5],
                              bullets=meta[6], screenshot_file=meta[7], side=meta[8])
        # чередуем: morph для smooth screenshot reveal, иногда fade
        if i % 3 == 0:
            add_transition(s, "morph", speed="med")
        elif i % 3 == 1:
            add_transition(s, "push", direction="l", speed="med")
        else:
            add_transition(s, "fade", speed="med")

    # === 17 Realtime ===
    s = slide_realtime(prs)
    add_transition(s, "split", speed="med")

    # === 18 Monitoring ===
    s = slide_monitoring(prs)
    add_transition(s, "push", direction="l", speed="med")

    # === 19 Testing ===
    s = slide_testing(prs)
    add_transition(s, "wipe", direction="l", speed="med")

    # === 20 Final ===
    s = slide_final(prs)
    add_transition(s, "fade", speed="slow")

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)

    size_kb = OUT_FILE.stat().st_size / 1024
    print(f"OK · {OUT_FILE.name} · {size_kb:.0f} KB · {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
